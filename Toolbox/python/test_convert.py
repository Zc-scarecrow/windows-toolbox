"""Quick integration test for Python converters."""
import os
import sys

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from convert import convert_file

TEST_DIR = os.path.join(os.path.dirname(__file__), "test_output")


def ensure_dir():
    os.makedirs(TEST_DIR, exist_ok=True)


def create_docx(path):
    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("This is a paragraph with bold text.")
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Name"
    table.rows[0].cells[1].text = "Value"
    table.rows[1].cells[0].text = "Foo"
    table.rows[1].cells[1].text = "Bar"
    doc.save(path)


def create_xlsx(path):
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Name"
    ws["B1"] = "Value"
    ws["A2"] = "Foo"
    ws["B2"] = "Bar"
    wb.save(path)


def create_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_textbox(0, 0, 4000000, 2000000).text_frame.text = "Hello PPTX"
    prs.save(path)


def main():
    ensure_dir()
    cases = []

    docx_path = os.path.join(TEST_DIR, "sample.docx")
    create_docx(docx_path)
    cases.append((docx_path, "docx", "md", "docx_to_md.md"))
    cases.append((docx_path, "docx", "txt", "docx_to_txt.txt"))
    cases.append((docx_path, "docx", "html", "docx_to_html.html"))

    xlsx_path = os.path.join(TEST_DIR, "sample.xlsx")
    create_xlsx(xlsx_path)
    cases.append((xlsx_path, "xlsx", "md", "xlsx_to_md.md"))
    cases.append((xlsx_path, "xlsx", "csv", "xlsx_to_csv.csv"))
    cases.append((xlsx_path, "xlsx", "json", "xlsx_to_json.json"))

    pptx_path = os.path.join(TEST_DIR, "sample.pptx")
    create_pptx(pptx_path)
    cases.append((pptx_path, "pptx", "md", "pptx_to_md.md"))
    cases.append((pptx_path, "pptx", "txt", "pptx_to_txt.txt"))

    failed = []
    for input_path, source, target, out_name in cases:
        output_path = os.path.join(TEST_DIR, out_name)
        print(f"Testing {source} -> {target}: {out_name}")
        try:
            convert_file(input_path, source, target, output_path)
            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                raise RuntimeError("Output file is empty or missing")
            print(f"  OK ({os.path.getsize(output_path)} bytes)")
        except Exception as e:
            print(f"  FAILED: {e}")
            failed.append((source, target, str(e)))

    print()
    if failed:
        print(f"FAILED: {len(failed)} case(s)")
        for s, t, e in failed:
            print(f"  {s} -> {t}: {e}")
        sys.exit(1)
    else:
        print("All tests passed.")


if __name__ == "__main__":
    main()
