"""PPTX conversion handler."""


class PptxConverter:
    """Convert PPTX to Markdown, HTML, and text."""

    def convert(self, input_path, target_format, output_path):
        if target_format == "md":
            self._to_markdown(input_path, output_path)
        elif target_format == "html":
            self._to_html(input_path, output_path)
        elif target_format == "txt":
            self._to_text(input_path, output_path)
        else:
            raise ValueError(f"Unsupported PPTX target format: {target_format}")

    def _extract_shape_text(self, shape, depth=0):
        """Recursively extract text from a shape (including group shapes)."""
        texts = []
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                child_texts = self._extract_shape_text(child, depth + 1)
                if child_texts:
                    texts.extend(child_texts)
        return texts

    def _to_markdown(self, input_path, output_path):
        from pptx import Presentation
        prs = Presentation(input_path)
        parts = []

        for idx, slide in enumerate(prs.slides, start=1):
            parts.append(f"## Slide {idx}")
            parts.append("")
            for shape in slide.shapes:
                texts = self._extract_shape_text(shape)
                if texts:
                    parts.append("\n\n".join(texts))
                    parts.append("")
            parts.append("")

        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(parts))

    def _to_html(self, input_path, output_path):
        from pptx import Presentation
        prs = Presentation(input_path)
        html_parts = [
            "<!DOCTYPE html>",
            "<html><head><meta charset=\"utf-8\"></head><body>",
        ]

        for idx, slide in enumerate(prs.slides, start=1):
            html_parts.append(f"<h2>Slide {idx}</h2>")
            html_parts.append("<div style=\"margin-bottom:24px\">")
            for shape in slide.shapes:
                texts = self._extract_shape_text(shape)
                if texts:
                    for t in texts:
                        escaped = t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                        html_parts.append(f"<p>{escaped}</p>")
            html_parts.append("</div>")

        html_parts.append("</body></html>")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(html_parts))

    def _to_text(self, input_path, output_path):
        from pptx import Presentation
        prs = Presentation(input_path)
        lines = []
        for idx, slide in enumerate(prs.slides, start=1):
            lines.append(f"--- Slide {idx} ---")
            for shape in slide.shapes:
                texts = self._extract_shape_text(shape)
                if texts:
                    lines.extend(texts)
            lines.append("")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
